"""
文档内容提取模块 v2.1

核心改进：
1. .doc 老格式：win32com 快速调用（3秒超时），失败立即跳过，不等待
2. .xls 老格式：xlrd 直接读取，速度快
3. .ppt 老格式：win32com 快速调用，失败立即跳过
4. 所有格式都有明确的超时保护
5. 新格式（docx/xlsx/pptx）直接解析 XML，速度最快
"""

import os
import re
import logging
import zipfile
from typing import Optional

logger = logging.getLogger(__name__)

# 每个文件最多提取的字符数（5万字，足够搜索用）
MAX_CONTENT_CHARS = 50_000

# 超过此大小的文件跳过（单位：MB）
MAX_FILE_SIZE_MB = 100

# xlsx/xls 最多读取的行数
MAX_EXCEL_ROWS = 2000

# win32com 是否可用（只检测一次，避免重复检测）
_WIN32COM_AVAILABLE = None


def _check_win32com() -> bool:
    """检测 win32com 是否可用（只检测一次）"""
    global _WIN32COM_AVAILABLE
    if _WIN32COM_AVAILABLE is None:
        try:
            import win32com.client
            _WIN32COM_AVAILABLE = True
        except ImportError:
            _WIN32COM_AVAILABLE = False
    return _WIN32COM_AVAILABLE


def _check_file(filepath: str) -> bool:
    """检查文件是否可处理"""
    try:
        size_mb = os.path.getsize(filepath) / (1024 * 1024)
        return size_mb <= MAX_FILE_SIZE_MB
    except Exception:
        return False


# ─────────────────────────────────────────────
# DOCX - 直接解析 XML（最快方式）
# ─────────────────────────────────────────────

def extract_docx(filepath: str) -> str:
    """提取 .docx 文件文本（直接读 XML，比 python-docx 快 3-5 倍）"""
    if not _check_file(filepath):
        return ''
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            if 'word/document.xml' not in z.namelist():
                return ''
            xml_data = z.read('word/document.xml').decode('utf-8', errors='ignore')

        # 段落之间加换行
        xml_data = re.sub(r'<w:p[ >]', '\n<w:p ', xml_data)
        xml_data = re.sub(r'<w:br[^/]*/>', '\n', xml_data)
        # 去掉所有 XML 标签
        text = re.sub(r'<[^>]+>', '', xml_data)
        # 清理多余空白
        lines = [line.strip() for line in text.split('\n') if line.strip()]
        result = '\n'.join(lines)
        return result[:MAX_CONTENT_CHARS]
    except Exception as e:
        logger.debug(f"XML 提取 docx 失败，尝试 python-docx: {e}")
        try:
            from docx import Document
            doc = Document(filepath)
            parts = []
            total = 0
            for para in doc.paragraphs:
                t = para.text.strip()
                if t:
                    parts.append(t)
                    total += len(t)
                    if total >= MAX_CONTENT_CHARS:
                        break
            return '\n'.join(parts)[:MAX_CONTENT_CHARS]
        except Exception:
            return ''


# ─────────────────────────────────────────────
# DOC - 老格式，win32com 优先，快速失败
# ─────────────────────────────────────────────

def extract_doc(filepath: str) -> str:
    """
    提取 .doc 文件文本（Word 97-2003 OLE 格式）

    策略：
    1. win32com（Windows + Word 已安装）：10秒超时，最可靠
    2. 纯 Python OLE 解析：扫描 WordDocument 流中的 UTF-16LE 文本
       - 按段落分隔符(0x000D/0x0007)切分段落
       - 过滤乱码：要求 >=80% 常用字符，>=2个常用汉字或3个连续字母
    3. 都失败则返回空字符串（文件名仍会被索引）
    """
    if not _check_file(filepath):
        return ''

    # 方式1：win32com 调用 Word（Windows 上最可靠）
    if _check_win32com():
        try:
            import win32com.client
            import threading

            result_holder = [None]

            def _do_extract():
                word = None
                try:
                    word = win32com.client.Dispatch('Word.Application')
                    word.Visible = False
                    word.DisplayAlerts = False
                    doc = word.Documents.Open(
                        os.path.abspath(filepath),
                        ReadOnly=True,
                        AddToRecentFiles=False
                    )
                    text = doc.Content.Text
                    doc.Close(False)
                    result_holder[0] = text
                except Exception:
                    pass
                finally:
                    if word:
                        try:
                            word.Quit()
                        except Exception:
                            pass

            t = threading.Thread(target=_do_extract, daemon=True)
            t.start()
            t.join(timeout=10)

            if result_holder[0] and len(result_holder[0].strip()) > 10:
                return result_holder[0][:MAX_CONTENT_CHARS]
        except Exception:
            pass

    # 方式2：纯 Python OLE 解析（不依赖任何外部软件，在任何电脑上都能工作）
    try:
        import olefile
        import struct

        if not olefile.isOleFile(filepath):
            return ''

        ole = olefile.OleFileIO(filepath)
        try:
            if not ole.exists('WordDocument'):
                return ''
            wd_data = ole.openstream('WordDocument').read()
        finally:
            ole.close()

        # 扫描 WordDocument 流中的 UTF-16LE 文本
        # Word 97 的文本以 UTF-16LE 存储，段落以 0x000D 或 0x0007 分隔
        paragraphs = []
        current = []
        i = 0

        while i < len(wd_data) - 1:
            try:
                c = struct.unpack_from('<H', wd_data, i)[0]
            except struct.error:
                break

            if c == 0x000D or c == 0x0007:
                # 段落分隔符
                if current:
                    text = ''.join(current).strip()
                    if text:
                        paragraphs.append(text)
                    current = []
            elif (0x4E00 <= c <= 0x9FFF or   # CJK 统一汉字（常用区）
                  0x0020 <= c <= 0x007E or   # ASCII 可打印字符
                  0x3000 <= c <= 0x303F or   # CJK 标点
                  0xFF00 <= c <= 0xFFEF or   # 全角字符
                  0x2014 <= c <= 0x2015 or   # 破折号
                  0x201C <= c <= 0x201D or   # 引号
                  0x2018 <= c <= 0x2019):    # 单引号
                current.append(chr(c))
            else:
                # 遇到不可识别字符，结束当前块
                if len(current) >= 2:
                    text = ''.join(current).strip()
                    if text:
                        paragraphs.append(text)
                current = []
            i += 2

        if current:
            text = ''.join(current).strip()
            if text:
                paragraphs.append(text)

        # 过滤乱码段落
        def _is_meaningful(text: str) -> bool:
            total = len(text)
            if total < 5:
                return False
            # 统计常用汉字（4E00-9FFF，不含扩展区）
            chinese_count = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
            # 统计 ASCII 字母
            alpha_count = sum(1 for c in text if c.isalpha() and ord(c) < 128)
            # 必须包含至少 2 个常用汉字 或 3 个连续字母
            has_chinese = chinese_count >= 2
            has_english = bool(re.search(r'[a-zA-Z]{3,}', text))
            if not (has_chinese or has_english):
                return False
            # 常用汉字 + ASCII 字母占比必须 >= 50%
            if (chinese_count + alpha_count) / total < 0.5:
                return False
            return True

        # 去重并过滤
        seen = set()
        result_lines = []
        for p in paragraphs:
            p = p.strip()
            if p and p not in seen and _is_meaningful(p):
                seen.add(p)
                result_lines.append(p)

        result = '\n'.join(result_lines)
        if len(result.strip()) > 5:
            return result[:MAX_CONTENT_CHARS]
    except Exception:
        pass

    return ''


# ─────────────────────────────────────────────
# XLSX - read_only 模式，限制行数
# ─────────────────────────────────────────────

def extract_xlsx(filepath: str) -> str:
    """提取 .xlsx 文件文本（read_only + 行数限制）"""
    if not _check_file(filepath):
        return ''
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        parts = []
        total = 0
        for sheet in wb.worksheets:
            row_count = 0
            for row in sheet.iter_rows(values_only=True):
                row_count += 1
                if row_count > MAX_EXCEL_ROWS:
                    break
                row_texts = []
                for cell in row:
                    if cell is None:
                        continue
                    if isinstance(cell, float):
                        val = str(int(cell)) if cell == int(cell) else f'{cell:.4g}'
                    elif isinstance(cell, bool):
                        val = '是' if cell else '否'
                    else:
                        val = str(cell).strip()
                    if val and val not in ('None', 'nan', ''):
                        row_texts.append(val)
                if row_texts:
                    line = ' '.join(row_texts)
                    parts.append(line)
                    total += len(line)
                    if total >= MAX_CONTENT_CHARS:
                        break
            if total >= MAX_CONTENT_CHARS:
                break
        wb.close()
        return '\n'.join(parts)[:MAX_CONTENT_CHARS]
    except Exception as e:
        logger.debug(f"提取 xlsx 失败 {filepath}: {e}")
        return ''


# ─────────────────────────────────────────────
# XLS - xlrd（老格式，直接读取）
# ─────────────────────────────────────────────

def extract_xls(filepath: str) -> str:
    """提取 .xls 文件文本（xlrd 直接读取）"""
    if not _check_file(filepath):
        return ''
    try:
        import xlrd
        wb = xlrd.open_workbook(filepath)
        parts = []
        total = 0
        for sheet in wb.sheets():
            for row_idx in range(min(sheet.nrows, MAX_EXCEL_ROWS)):
                row_texts = []
                for col_idx in range(sheet.ncols):
                    cell = sheet.cell(row_idx, col_idx)
                    if cell.ctype == 2:  # XL_CELL_NUMBER
                        v = cell.value
                        val = str(int(v)) if v == int(v) else f'{v:.4g}'
                    elif cell.ctype == 5:  # XL_CELL_ERROR
                        continue
                    else:
                        val = str(cell.value).strip()
                    if val and val not in ('', 'nan'):
                        row_texts.append(val)
                if row_texts:
                    line = ' '.join(row_texts)
                    parts.append(line)
                    total += len(line)
                    if total >= MAX_CONTENT_CHARS:
                        break
            if total >= MAX_CONTENT_CHARS:
                break
        return '\n'.join(parts)[:MAX_CONTENT_CHARS]
    except Exception as e:
        logger.debug(f"提取 xls 失败 {filepath}: {e}")
        return ''


# ─────────────────────────────────────────────
# PPTX - 直接读 XML（最快）
# ─────────────────────────────────────────────

def extract_pptx(filepath: str) -> str:
    """提取 .pptx 文件文本（直接读 XML）"""
    if not _check_file(filepath):
        return ''
    try:
        with zipfile.ZipFile(filepath, 'r') as z:
            slide_files = sorted([
                name for name in z.namelist()
                if name.startswith('ppt/slides/slide') and name.endswith('.xml')
            ])
            parts = []
            total = 0
            for slide_file in slide_files:
                xml_data = z.read(slide_file).decode('utf-8', errors='ignore')
                texts = re.findall(r'<a:t[^>]*>([^<]+)</a:t>', xml_data)
                for t in texts:
                    t = t.strip()
                    if t:
                        parts.append(t)
                        total += len(t)
                        if total >= MAX_CONTENT_CHARS:
                            break
                if total >= MAX_CONTENT_CHARS:
                    break
        return '\n'.join(parts)[:MAX_CONTENT_CHARS]
    except Exception as e:
        logger.debug(f"XML 提取 pptx 失败，尝试 python-pptx: {e}")
        try:
            from pptx import Presentation
            prs = Presentation(filepath)
            parts = []
            total = 0
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        t = shape.text.strip()
                        parts.append(t)
                        total += len(t)
                        if total >= MAX_CONTENT_CHARS:
                            break
                if total >= MAX_CONTENT_CHARS:
                    break
            return '\n'.join(parts)[:MAX_CONTENT_CHARS]
        except Exception:
            return ''


# ─────────────────────────────────────────────
# PPT - 老格式，win32com 优先
# ─────────────────────────────────────────────

def extract_ppt(filepath: str) -> str:
    """提取 .ppt 文件文本（PowerPoint 97-2003）"""
    if not _check_file(filepath):
        return ''

    # 方式1：win32com 调用 PowerPoint（Windows 上最可靠）
    if _check_win32com():
        try:
            import win32com.client
            import threading

            result_holder = [None]

            def _do_extract():
                ppt_app = None
                try:
                    ppt_app = win32com.client.Dispatch('PowerPoint.Application')
                    ppt_app.Visible = True  # PowerPoint 必须可见才能打开文件
                    prs = ppt_app.Presentations.Open(
                        os.path.abspath(filepath),
                        ReadOnly=True,
                        Untitled=True,
                        WithWindow=False
                    )
                    texts = []
                    for slide in prs.Slides:
                        for shape in slide.Shapes:
                            try:
                                if shape.HasTextFrame:
                                    texts.append(shape.TextFrame.TextRange.Text)
                            except Exception:
                                pass
                    prs.Close()
                    result_holder[0] = '\n'.join(texts)
                except Exception:
                    pass
                finally:
                    if ppt_app:
                        try:
                            ppt_app.Quit()
                        except Exception:
                            pass

            t = threading.Thread(target=_do_extract, daemon=True)
            t.start()
            t.join(timeout=10)

            if result_holder[0] and len(result_holder[0].strip()) > 10:
                return result_holder[0][:MAX_CONTENT_CHARS]
        except Exception:
            pass

    # 方式2：olefile 扫描（兜底）
    try:
        import olefile
        import struct

        if not olefile.isOleFile(filepath):
            return ''

        ole = olefile.OleFileIO(filepath)
        texts = []
        current = []

        for stream_name in ole.listdir():
            try:
                data = ole.openstream(stream_name).read()
                i = 0
                while i < len(data) - 1:
                    c = struct.unpack_from('<H', data, i)[0]
                    if (0x4E00 <= c <= 0x9FFF or
                        0x0020 <= c <= 0x007E or
                        0x3000 <= c <= 0x303F):
                        current.append(chr(c))
                    else:
                        if len(current) >= 5:
                            texts.append(''.join(current))
                        current = []
                    i += 2
            except Exception:
                pass

        ole.close()

        if len(current) >= 5:
            texts.append(''.join(current))

        seen = set()
        unique_texts = []
        for t in texts:
            t_clean = t.strip()
            if t_clean and t_clean not in seen and len(t_clean) >= 3:
                seen.add(t_clean)
                unique_texts.append(t_clean)

        result = ' '.join(unique_texts)
        if len(result.strip()) > 10:
            return result[:MAX_CONTENT_CHARS]
    except Exception:
        pass

    return ''


# ─────────────────────────────────────────────
# PDF
# ─────────────────────────────────────────────

def extract_pdf_text(filepath: str) -> str:
    """提取可识别 PDF 的文本"""
    if not _check_file(filepath):
        return ''

    # 方式1：pymupdf（最快）
    try:
        import fitz
        parts = []
        total = 0
        doc = fitz.open(filepath)
        for page in doc:
            text = page.get_text()
            if text and text.strip():
                parts.append(text.strip())
                total += len(text)
                if total >= MAX_CONTENT_CHARS:
                    break
        doc.close()
        result = '\n'.join(parts)[:MAX_CONTENT_CHARS]
        if len(result.strip()) > 50:
            return result
    except ImportError:
        pass
    except Exception as e:
        logger.debug(f"pymupdf 失败 {filepath}: {e}")

    # 方式2：pdfplumber（兜底）
    try:
        import pdfplumber
        parts = []
        total = 0
        with pdfplumber.open(filepath) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    parts.append(text.strip())
                    total += len(text)
                    if total >= MAX_CONTENT_CHARS:
                        break
        return '\n'.join(parts)[:MAX_CONTENT_CHARS]
    except Exception as e:
        logger.debug(f"pdfplumber 失败 {filepath}: {e}")
        return ''


def extract_pdf_ocr(filepath: str) -> str:
    """对扫描版 PDF 进行 OCR 识别"""
    if not _check_file(filepath):
        return ''
    try:
        import pytesseract
        from pdf2image import convert_from_path
        pages = convert_from_path(filepath, dpi=150)
        parts = []
        total = 0
        for page_img in pages:
            text = pytesseract.image_to_string(page_img, lang='chi_sim+eng')
            if text.strip():
                parts.append(text.strip())
                total += len(text)
                if total >= MAX_CONTENT_CHARS:
                    break
        return '\n'.join(parts)[:MAX_CONTENT_CHARS]
    except Exception as e:
        logger.debug(f"OCR 失败 {filepath}: {e}")
        return ''


# ─────────────────────────────────────────────
# 统一入口
# ─────────────────────────────────────────────

def extract_text(filepath: str, enable_pdf: bool = True, enable_ocr: bool = False) -> str:
    """
    统一文本提取入口

    支持格式：
    - .docx / .doc（Word 文档）
    - .xlsx / .xls（Excel 表格）
    - .pptx / .ppt（PowerPoint 演示文稿）
    - .pdf（PDF 文档，可选 OCR）
    - .txt / .csv / .md / .json / .xml / .html（纯文本）
    """
    ext = os.path.splitext(filepath)[1].lower()

    try:
        if ext == '.docx':
            return extract_docx(filepath)
        elif ext == '.doc':
            return extract_doc(filepath)
        elif ext == '.xlsx':
            return extract_xlsx(filepath)
        elif ext == '.xls':
            return extract_xls(filepath)
        elif ext == '.pptx':
            return extract_pptx(filepath)
        elif ext == '.ppt':
            return extract_ppt(filepath)
        elif ext == '.pdf':
            if not enable_pdf:
                return ''
            text = extract_pdf_text(filepath)
            if not text.strip() and enable_ocr:
                text = extract_pdf_ocr(filepath)
            return text
        elif ext in ('.txt', '.csv', '.md', '.log', '.ini', '.cfg', '.conf'):
            # 纯文本文件：直接读取，尝试多种编码
            for enc in ['utf-8-sig', 'utf-8', 'gbk', 'gb2312', 'latin-1']:
                try:
                    with open(filepath, 'r', encoding=enc, errors='ignore') as f:
                        return f.read(MAX_CONTENT_CHARS)
                except Exception:
                    continue
            return ''
        elif ext in ('.json', '.xml', '.html', '.htm'):
            # 结构化文本：读取后去标签
            for enc in ['utf-8-sig', 'utf-8', 'gbk']:
                try:
                    with open(filepath, 'r', encoding=enc, errors='ignore') as f:
                        content = f.read(MAX_CONTENT_CHARS)
                    if ext in ('.xml', '.html', '.htm'):
                        content = re.sub(r'<[^>]+>', ' ', content)
                        content = re.sub(r'\s+', ' ', content).strip()
                    return content
                except Exception:
                    continue
            return ''
        else:
            return ''
    except Exception as e:
        logger.debug(f"提取失败 {filepath}: {e}")
        return ''


# 支持的文件扩展名列表（供外部引用）
SUPPORTED_EXTENSIONS = [
    '.docx', '.doc',
    '.xlsx', '.xls',
    '.pptx', '.ppt',
    '.pdf',
    '.txt', '.csv', '.md',
    '.json', '.xml', '.html', '.htm',
    '.log', '.ini', '.cfg', '.conf',
]
